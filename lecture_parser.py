from pptx import Presentation, slide


class LectureParser:
    def __init__(self, lecture_path: str):
        self._lecture = lecture_path

    def get_lecture_parts(self):
        raise NotImplementedError

    def parse_lecture_part(self, lecture_part) -> str:
        raise NotImplementedError


class PresentationParser(LectureParser):
    def __init__(self, lecture_path):
        super().__init__(lecture_path)
        self._parser = Presentation(self._lecture)

    def get_lecture_parts(self):
        for pres_slide in self._parser.slides:
            yield pres_slide

    def parse_lecture_part(self, pres_slide: slide.Slide) -> str:
        # text_runs will be populated with a list of strings,
        # one for each text run in presentation
        text_runs = []
        for shape in pres_slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text_runs.append(run.text)
        return " ".join(text_runs)


def lecture_factory(lecture_type: str) -> type[LectureParser]:
    lecture_type_dict = {
        "pptx": PresentationParser
        # "pdf": PdfParser
    }
    return lecture_type_dict[lecture_type]

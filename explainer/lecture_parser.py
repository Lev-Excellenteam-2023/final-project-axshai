from pptx import Presentation, slide


class LectureParser:
    def __init__(self, lecture_path: str):
        """
        Initializes a LectureParser object with the path to the lecture file.

        :param lecture_path: The path to the lecture file.
        """
        self._lecture = lecture_path

    def get_lecture_parts(self):
        """
        Retrieves the lecture parts.

        This method is meant to be overridden by subclasses.

        :return: An iterable of lecture parts.
        """
        raise NotImplementedError

    def parse_lecture_part(self, lecture_part) -> str:
        """
        Parses a lecture part and returns the extracted text.

        This method is meant to be overridden by subclasses.

        :param lecture_part: The lecture part to parse.
        :return: The extracted text from the lecture part.
        """
        raise NotImplementedError


class PresentationParser(LectureParser):
    def __init__(self, lecture_path):
        """
        Initializes a PresentationParser object with the path to the PowerPoint presentation.

        :param lecture_path: The path to the PowerPoint presentation.
        """
        super().__init__(lecture_path)
        self._parser = Presentation(self._lecture)

    def get_lecture_parts(self):
        """
        Retrieves the slides from the PowerPoint presentation.

        :return: An iterator of slides in the presentation.
        """
        for pres_slide in self._parser.slides:
            yield pres_slide

    def parse_lecture_part(self, pres_slide: slide.Slide) -> str:
        """
        Parses a slide from the PowerPoint presentation and returns the extracted text.

        :param pres_slide: The slide to parse.
        :return: The extracted text from the slide.
        """
        text_runs = [' '.join(run.text.split()) for shape in pres_slide.shapes if shape.has_text_frame for paragraph in
                     shape.text_frame.paragraphs for run in paragraph.runs]
        return " ".join(text_runs)


def lecture_factory(lecture_type: str) -> type[LectureParser]:
    """
    Returns a LectureParser subclass based on the provided lecture type.

    :param lecture_type: The type of the lecture (e.g., "pptx", "pdf").
    :return: The corresponding LectureParser subclass.
    """
    lecture_type_dict = {
        "pptx": PresentationParser
        # "pdf": PdfParser
    }
    return lecture_type_dict[lecture_type]
